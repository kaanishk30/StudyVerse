from flask import Flask, render_template, request, jsonify, session, redirect, url_for, Response
import numpy as np
import joblib
import requests
from bs4 import BeautifulSoup
import wikipediaapi
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
import re
import os
from werkzeug.utils import secure_filename
from werkzeug.security import generate_password_hash, check_password_hash
import PyPDF2
from pptx import Presentation
import random
from datetime import datetime, timedelta
import json
import sqlite3
from functools import wraps
import time
import uuid

# ==========================================
# FLASK APP INITIALIZATION
# ==========================================
app = Flask(__name__)
app.secret_key = 'change-this-to-a-random-secret-key-in-production'

UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'pdf', 'ppt', 'pptx', 'txt'}
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs('models', exist_ok=True)

# ==========================================
# DATABASE SETUP
# ==========================================
def init_db():
    """Initialize SQLite database"""
    conn = sqlite3.connect('studypal.db')
    c = conn.cursor()
    
    c.execute('''
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT UNIQUE NOT NULL,
            email TEXT UNIQUE NOT NULL,
            password_hash TEXT NOT NULL,
            full_name TEXT NOT NULL,
            bio TEXT,
            avatar_url TEXT,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            last_login TIMESTAMP
        )
    ''')
    
    c.execute('''
        CREATE TABLE IF NOT EXISTS user_stats (
            user_id INTEGER PRIMARY KEY,
            streak INTEGER DEFAULT 0,
            last_activity DATE,
            total_quizzes INTEGER DEFAULT 0,
            correct_answers INTEGER DEFAULT 0,
            total_questions INTEGER DEFAULT 0,
            total_study_time INTEGER DEFAULT 0,
            FOREIGN KEY (user_id) REFERENCES users(id)
        )
    ''')
    
    c.execute('''
        CREATE TABLE IF NOT EXISTS study_history (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER,
            topic TEXT,
            score TEXT,
            difficulty TEXT,
            study_duration INTEGER DEFAULT 0,
            quiz_data TEXT,
            timestamp TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (user_id) REFERENCES users(id)
        )
    ''')
    
    c.execute('''
        CREATE TABLE IF NOT EXISTS user_achievements (
            user_id INTEGER,
            achievement_id TEXT,
            unlocked_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            PRIMARY KEY (user_id, achievement_id),
            FOREIGN KEY (user_id) REFERENCES users(id)
        )
    ''')
    
    c.execute('''
        CREATE TABLE IF NOT EXISTS schedules (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER,
            subject TEXT NOT NULL,
            unit_name TEXT,
            topic_name TEXT NOT NULL,
            difficulty TEXT DEFAULT 'medium',
            estimated_hours INTEGER DEFAULT 2,
            scheduled_date DATE NOT NULL,
            start_time TEXT,
            end_time TEXT,
            status TEXT DEFAULT 'pending',
            completion_percentage INTEGER DEFAULT 0,
            notes TEXT,
            is_auto_generated INTEGER DEFAULT 0,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            completed_at TIMESTAMP,
            FOREIGN KEY (user_id) REFERENCES users(id)
        )
    ''')
    
    # Temporary table to store quiz sessions (avoids session cookie size limits)
    c.execute('''
        CREATE TABLE IF NOT EXISTS temp_quiz_sessions (
            session_id TEXT PRIMARY KEY,
            user_id INTEGER,
            quiz_data TEXT NOT NULL,
            topic TEXT,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (user_id) REFERENCES users(id)
        )
    ''')
    
    # Clean up old quiz sessions (older than 24 hours)
    c.execute('''
        DELETE FROM temp_quiz_sessions 
        WHERE created_at < datetime('now', '-1 day')
    ''')
    
    conn.commit()
    conn.close()
    print("✓ Database initialized")

init_db()

# ==========================================
# AI SCHEDULER ENGINE
# ==========================================
class AIScheduler:
    """AI-powered schedule generator with difficulty estimation"""
    
    DIFFICULTY_KEYWORDS = {
        'easy': ['introduction', 'basic', 'overview', 'fundamentals', 'simple', 'getting started'],
        'medium': ['intermediate', 'application', 'practice', 'implementation', 'working with'],
        'hard': ['advanced', 'complex', 'optimization', 'theory', 'deep dive', 'architecture']
    }
    
    STUDY_TIME_BY_DIFFICULTY = {
        'easy': 2,      # 2 hours
        'medium': 3,    # 3 hours
        'hard': 4       # 4 hours
    }
    
    @staticmethod
    def parse_syllabus(syllabus_text):
        """Parse syllabus text and extract units/topics"""
        lines = [line.strip() for line in syllabus_text.split('\n') if line.strip()]
        structure = []
        current_unit = None
        
        for line in lines:
            # Detect unit headers (Unit 1, Module 1, Chapter 1, etc.)
            unit_match = re.match(r'^(Unit|Module|Chapter|Week)\s*(\d+)[:\s-]*(.*)$', line, re.IGNORECASE)
            
            if unit_match:
                unit_type = unit_match.group(1)
                unit_num = unit_match.group(2)
                unit_name = unit_match.group(3).strip() or f"{unit_type} {unit_num}"
                current_unit = {
                    'unit_name': unit_name,
                    'topics': []
                }
                structure.append(current_unit)
            elif current_unit is not None:
                # Add as topic to current unit
                # Clean topic name - remove leading bullets/dashes
                topic = re.sub(r'^[-*•]\s*', '', line)
                if len(topic) > 3:  # Ignore very short lines
                    current_unit['topics'].append(topic)
            else:
                # No unit detected yet, treat as standalone topic
                # Clean topic name
                topic = re.sub(r'^[-*•]\s*', '', line)
                if len(topic) > 3:
                    if not structure:
                        structure.append({
                            'unit_name': 'General Topics',
                            'topics': []
                        })
                    structure[0]['topics'].append(topic)
        
        return structure
    
    @staticmethod
    def estimate_difficulty(topic_name):
        """Estimate topic difficulty using balanced keyword matching and complexity analysis"""
        if not topic_name or len(topic_name.strip()) < 3:
            return 'medium'
        
        topic_lower = topic_name.lower()
        topic_words = topic_name.split()
        
        # Enhanced difficulty keywords with more comprehensive coverage
        easy_keywords = [
            'introduction', 'intro', 'basic', 'basics', 'overview', 'fundamentals', 
            'simple', 'getting started', 'beginner', 'elementary', 'foundation',
            'what is', 'definition', 'concept', 'understanding', 'primer',
            'first', 'starting', 'beginning', 'initial', 'essentials'
        ]
        
        medium_keywords = [
            'intermediate', 'application', 'practice', 'implementation', 'working with',
            'using', 'applying', 'methods', 'techniques', 'process', 'procedure',
            'analysis', 'comparison', 'types', 'classification', 'structure',
            'design', 'development', 'building', 'creating', 'programming'
        ]
        
        hard_keywords = [
            'advanced', 'complex', 'optimization', 'theory', 'deep dive', 'architecture',
            'algorithm', 'mathematical', 'proof', 'derivation', 'research', 'cutting-edge',
            'sophisticated', 'intricate', 'comprehensive', 'in-depth', 'expert',
            'quantum', 'neural', 'distributed', 'concurrent', 'cryptographic',
            'theoretical', 'abstract', 'formal'
        ]
        
        # Count keyword matches (case-insensitive, whole word matching)
        easy_score = 0
        medium_score = 0
        hard_score = 0
        
        for keyword in easy_keywords:
            if keyword in topic_lower:
                easy_score += 2  # Strong indicator
        
        for keyword in medium_keywords:
            if keyword in topic_lower:
                medium_score += 2
        
        for keyword in hard_keywords:
            if keyword in topic_lower:
                hard_score += 2
        
        # If clear keyword match found, use it
        if easy_score > medium_score and easy_score > hard_score and easy_score > 0:
            return 'easy'
        if hard_score > easy_score and hard_score > medium_score and hard_score > 0:
            return 'hard'
        if medium_score > 0:
            return 'medium'
        
        # Complexity analysis based on topic characteristics
        complexity_score = 0
        
        # 1. Word count analysis (shorter topics tend to be simpler)
        word_count = len(topic_words)
        if word_count <= 2:
            complexity_score += 0  # Very short, likely easy
        elif word_count <= 4:
            complexity_score += 1  # Short, likely easy-medium
        elif word_count <= 6:
            complexity_score += 2  # Medium length
        else:
            complexity_score += 3  # Long, likely complex
        
        # 2. Technical term detection (CamelCase, numbers, special chars)
        technical_terms = 0
        for word in topic_words:
            # Check for CamelCase (e.g., "JavaScript", "TensorFlow")
            if len(word) > 1 and any(c.isupper() for c in word[1:]):
                technical_terms += 1
            # Check for numbers (e.g., "Python3", "v2.0")
            if any(c.isdigit() for c in word):
                technical_terms += 1
        
        if technical_terms == 0:
            complexity_score += 0
        elif technical_terms == 1:
            complexity_score += 1
        else:
            complexity_score += 2
        
        # 3. Mathematical symbols or special characters
        has_math = any(char in topic_name for char in ['=', '+', '∫', '∑', '∂', '∆', 'α', 'β'])
        if has_math:
            complexity_score += 2
        
        # 4. Version numbers (usually indicates advanced/specific topics)
        has_version = bool(re.search(r'\d+\.\d+|\bv\d+\b', topic_lower))
        if has_version:
            complexity_score += 1
        
        # 5. Subject-specific difficulty patterns (known hard subjects)
        hard_subjects = [
            'calculus', 'quantum', 'thermodynamics', 'genetics', 'organic chemistry',
            'differential equations', 'linear algebra', 'topology', 'cryptography',
            'compiler', 'operating system', 'machine learning', 'deep learning'
        ]
        
        easy_subjects = [
            'html', 'css', 'basic', 'intro', 'overview', 'summary',
            'list', 'guide', 'tutorial', 'example'
        ]
        
        for subject in hard_subjects:
            if subject in topic_lower:
                complexity_score += 2
                break
        
        for subject in easy_subjects:
            if subject in topic_lower:
                complexity_score -= 2  # Reduce score for easy subjects
                break
        
        # 6. Check for action words that indicate difficulty
        if any(word in topic_lower for word in ['learn', 'understand', 'know', 'explore']):
            complexity_score -= 1  # Learning-focused, likely easier
        
        if any(word in topic_lower for word in ['master', 'expert', 'professional', 'advanced']):
            complexity_score += 2  # Mastery-focused, likely harder
        
        # Final difficulty assignment with balanced thresholds
        # Adjusted thresholds to prevent everything being marked as hard
        if complexity_score <= 2:
            return 'easy'
        elif complexity_score <= 5:
            return 'medium'
        else:
            return 'hard'
    
    @staticmethod
    def generate_schedule(syllabus_text, total_days, start_date=None):
        """Generate optimized study schedule"""
        try:
            if start_date is None:
                start_date = datetime.now().date()
            elif isinstance(start_date, str):
                try:
                    start_date = datetime.strptime(start_date, '%Y-%m-%d').date()
                except:
                    start_date = datetime.now().date()
            
            # Parse syllabus
            structure = AIScheduler.parse_syllabus(syllabus_text)
            
            if not structure:
                return []
            
            # Calculate total topics and estimated time
            all_topics = []
            for unit in structure:
                for topic in unit['topics']:
                    difficulty = AIScheduler.estimate_difficulty(topic)
                    estimated_hours = AIScheduler.STUDY_TIME_BY_DIFFICULTY[difficulty]
                    all_topics.append({
                        'unit_name': unit['unit_name'],
                        'topic_name': topic,
                        'difficulty': difficulty,
                        'estimated_hours': estimated_hours
                    })
            
            if not all_topics:
                return []
            
            total_hours = sum(t['estimated_hours'] for t in all_topics)
            hours_per_day = total_hours / total_days if total_days > 0 else 3
            
            # Ensure reasonable hours per day (max 8 hours)
            if hours_per_day > 8:
                hours_per_day = 8
            
            # Generate schedule with time distribution
            schedule = []
            current_date = start_date
            current_day_hours = 0
            
            for topic in all_topics:
                # If adding this topic exceeds daily limit, move to next day
                if current_day_hours > 0 and current_day_hours + topic['estimated_hours'] > hours_per_day * 1.5:
                    current_date += timedelta(days=1)
                    current_day_hours = 0
                    
                    # Skip weekends if needed
                    while current_date.weekday() >= 5:  # 5=Saturday, 6=Sunday
                        current_date += timedelta(days=1)
                
                # Calculate time slots (morning, afternoon, evening)
                if current_day_hours == 0:
                    start_time = "09:00"
                elif current_day_hours < 4:
                    start_time = "14:00"
                else:
                    start_time = "18:00"
                
                # Calculate end time
                start_hour = int(start_time.split(':')[0])
                end_hour = start_hour + topic['estimated_hours']
                
                # Ensure end time doesn't exceed 22:00
                if end_hour > 22:
                    end_hour = 22
                
                end_time = f"{end_hour:02d}:00"
                
                schedule.append({
                    'unit_name': topic['unit_name'],
                    'topic_name': topic['topic_name'],
                    'difficulty': topic['difficulty'],
                    'estimated_hours': topic['estimated_hours'],
                    'scheduled_date': current_date.isoformat(),
                    'start_time': start_time,
                    'end_time': end_time,
                    'is_auto_generated': 1
                })
                
                current_day_hours += topic['estimated_hours']
            
            return schedule
        
        except Exception as e:
            print(f"Error in generate_schedule: {e}")
            import traceback
            traceback.print_exc()
            return []

# ==========================================
# AUTHENTICATION HELPERS
# ==========================================
def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'user_id' not in session:
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function

def get_db_connection():
    conn = sqlite3.connect('studypal.db')
    conn.row_factory = sqlite3.Row
    return conn

def get_user_by_id(user_id):
    conn = get_db_connection()
    user = conn.execute('SELECT * FROM users WHERE id = ?', (user_id,)).fetchone()
    conn.close()
    return dict(user) if user else None

def get_user_stats(user_id):
    conn = get_db_connection()
    stats = conn.execute('SELECT * FROM user_stats WHERE user_id = ?', (user_id,)).fetchone()
    conn.close()
    return dict(stats) if stats else None

def update_streak(user_id):
    """Update user streak with proper date handling"""
    conn = get_db_connection()
    stats = conn.execute('SELECT * FROM user_stats WHERE user_id = ?', (user_id,)).fetchone()
    
    today = datetime.now().date().isoformat()
    
    if stats:
        last_activity = stats['last_activity']
        current_streak = stats['streak']
        
        if last_activity:
            try:
                last_date = datetime.strptime(last_activity, '%Y-%m-%d').date()
                today_date = datetime.now().date()
                days_diff = (today_date - last_date).days
                
                if days_diff == 0:
                    # Same day, don't change streak
                    pass
                elif days_diff == 1:
                    # Consecutive day, increment
                    current_streak += 1
                    conn.execute('UPDATE user_stats SET streak = ?, last_activity = ? WHERE user_id = ?',
                                (current_streak, today, user_id))
                elif days_diff > 1:
                    # Streak broken, reset to 1
                    current_streak = 1
                    conn.execute('UPDATE user_stats SET streak = ?, last_activity = ? WHERE user_id = ?',
                                (current_streak, today, user_id))
            except:
                current_streak = 1
                conn.execute('UPDATE user_stats SET streak = ?, last_activity = ? WHERE user_id = ?',
                            (current_streak, today, user_id))
        else:
            current_streak = 1
            conn.execute('UPDATE user_stats SET streak = ?, last_activity = ? WHERE user_id = ?',
                        (current_streak, today, user_id))
    else:
        conn.execute('INSERT INTO user_stats (user_id, streak, last_activity) VALUES (?, 1, ?)',
                    (user_id, today))
        current_streak = 1
    
    conn.commit()
    conn.close()
    return current_streak

# ==========================================
# NLTK SETUP
# ==========================================
print("📦 Setting up NLTK...")
required_nltk = ['punkt', 'punkt_tab', 'stopwords', 'averaged_perceptron_tagger', 'averaged_perceptron_tagger_eng']
for dataset in required_nltk:
    try:
        nltk.download(dataset, quiet=True)
    except:
        pass
print("✓ NLTK ready")

# ==========================================
# MULTI-SOURCE CONTENT FETCHER
# ==========================================
class MultiSourceLearner:
    """Enhanced content fetcher with multiple sources including web search"""
    
    def __init__(self):
        self.wiki = wikipediaapi.Wikipedia(
            language='en',
            extract_format=wikipediaapi.ExtractFormat.WIKI,
            user_agent="AI_Study_Pal_Multi/1.0"
        )
    
    def clean_text(self, text):
        if not text:
            return ""
        # Remove wiki markup
        text = re.sub(r'={2,}[^=]*={2,}', '', text)
        text = re.sub(r'\[.*?\]', '', text)
        # Remove excessive newlines
        text = re.sub(r'\n{3,}', '\n\n', text)
        # Remove excessive spaces
        text = re.sub(r' {2,}', ' ', text)
        # Remove non-printable characters but keep basic punctuation
        text = re.sub(r'[^\x20-\x7E\n\r\t]', '', text)
        return text.strip()
    
    def fetch_from_wikipedia(self, query):
        """Source 1: Wikipedia"""
        try:
            print(f"  📖 [Wikipedia] Searching: {query}")
            search_terms = [query, query.title(), f"{query} (concept)", f"{query} (technology)"]
            
            for term in search_terms:
                page = self.wiki.page(term)
                if page.exists():
                    text = page.text if hasattr(page, 'text') else page.summary
                    if len(text) > 100:
                        print(f"  ✅ [Wikipedia] Found: {page.title}")
                        return self.clean_text(text[:5000]), f"Wikipedia: {page.title}"
        except Exception as e:
            print(f"  ⚠️ [Wikipedia] Error: {e}")
        return "", ""
    
    def fetch_from_simple_wikipedia(self, query):
        """Source 2: Simple English Wikipedia"""
        try:
            print(f"  📖 [Simple Wiki] Searching: {query}")
            simple_wiki = wikipediaapi.Wikipedia(
                language='simple',
                extract_format=wikipediaapi.ExtractFormat.WIKI,
                user_agent="AI_Study_Pal_Multi/1.0"
            )
            
            page = simple_wiki.page(query)
            if page.exists():
                text = page.text if hasattr(page, 'text') else page.summary
                if len(text) > 100:
                    print(f"  ✅ [Simple Wiki] Found: {page.title}")
                    return self.clean_text(text[:3000]), f"Simple Wikipedia: {page.title}"
        except Exception as e:
            print(f"  ⚠️ [Simple Wiki] Error: {e}")
        return "", ""
    
    def fetch_from_duckduckgo(self, query):
        """Source 3: DuckDuckGo Instant Answer"""
        try:
            print(f"  🔍 [DuckDuckGo] Searching: {query}")
            api_url = f"https://api.duckduckgo.com/?q={query}&format=json"
            
            response = requests.get(api_url, timeout=10)
            data = response.json()
            
            text = ""
            if data.get('AbstractText'):
                text = data['AbstractText']
            elif data.get('RelatedTopics'):
                for topic in data['RelatedTopics'][:5]:
                    if isinstance(topic, dict) and topic.get('Text'):
                        text += topic['Text'] + "\n\n"
            
            if len(text) > 100:
                print(f"  ✅ [DuckDuckGo] Found instant answer")
                return self.clean_text(text[:3000]), "DuckDuckGo Instant Answer"
        except Exception as e:
            print(f"  ⚠️ [DuckDuckGo] Error: {e}")
        return "", ""
    
    def fetch_from_web_search(self, query):
        """Source 4: Web scraping from search results"""
        try:
            print(f"  🌐 [Web Search] Searching: {query}")
            
            # Try to get content from educational sites
            search_query = query.replace(' ', '+')
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            }
            
            # Try Britannica
            try:
                britannica_url = f"https://www.britannica.com/search?query={search_query}"
                response = requests.get(britannica_url, headers=headers, timeout=10)
                if response.status_code == 200:
                    soup = BeautifulSoup(response.content, 'html.parser')
                    # Look for article content
                    content_divs = soup.find_all(['p', 'div'], class_=re.compile('content|article|text'))
                    text = ' '.join([div.get_text() for div in content_divs[:10]])
                    if len(text) > 200:
                        print(f"  ✅ [Britannica] Found content")
                        return self.clean_text(text[:4000]), "Encyclopedia Britannica"
            except:
                pass
            
        except Exception as e:
            print(f"  ⚠️ [Web Search] Error: {e}")
        return "", ""
    
    def fetch_comparison_content(self, query):
        """Source 5: Handle comparison queries (X vs Y)"""
        try:
            if ' vs ' in query.lower() or ' versus ' in query.lower():
                print(f"  🔄 [Comparison] Detected comparison query")
                
                # Split the query
                parts = re.split(r'\s+vs\.?\s+|\s+versus\s+', query, flags=re.IGNORECASE)
                if len(parts) == 2:
                    topic1, topic2 = parts[0].strip(), parts[1].strip()
                    
                    # Fetch both topics
                    content1, source1 = self.fetch_from_wikipedia(topic1)
                    time.sleep(0.5)
                    content2, source2 = self.fetch_from_wikipedia(topic2)
                    
                    if content1 and content2:
                        combined = f"COMPARISON: {topic1} vs {topic2}\n\n"
                        combined += f"=== {topic1.upper()} ===\n{content1[:2000]}\n\n"
                        combined += f"=== {topic2.upper()} ===\n{content2[:2000]}\n\n"
                        combined += f"KEY DIFFERENCES:\n"
                        combined += f"While {topic1} and {topic2} are related concepts, they have distinct characteristics. "
                        combined += f"{topic1} focuses on specific aspects, while {topic2} encompasses broader principles."
                        
                        print(f"  ✅ [Comparison] Generated comparison content")
                        return combined, f"Comparison: {source1} & {source2}"
        except Exception as e:
            print(f"  ⚠️ [Comparison] Error: {e}")
        return "", ""
    
    def fetch_from_fallback_sources(self, query):
        """Source 6: Fallback with general knowledge"""
        try:
            print(f"  📚 [Fallback] Generating general content for: {query}")
            
            # Create basic educational content structure
            content = f"TOPIC: {query}\n\n"
            content += f"This is an educational topic that covers various aspects of {query}. "
            content += f"To learn more about {query}, consider exploring the following areas:\n\n"
            content += f"1. Fundamental Concepts: Understanding the basic principles and definitions related to {query}.\n\n"
            content += f"2. Historical Context: How {query} developed over time and its evolution.\n\n"
            content += f"3. Practical Applications: Real-world uses and implementations of {query}.\n\n"
            content += f"4. Key Components: The main elements that make up {query}.\n\n"
            content += f"5. Related Topics: Other subjects that connect to {query}.\n\n"
            content += f"For more detailed information, try:\n"
            content += f"- Searching for specific aspects of {query}\n"
            content += f"- Uploading your own study materials\n"
            content += f"- Adding your class notes in the text area\n"
            
            print(f"  ✅ [Fallback] Generated basic content structure")
            return content, "AI-Generated Study Guide"
            
        except Exception as e:
            print(f"  ⚠️ [Fallback] Error: {e}")
        return "", ""
    
    def search_and_learn(self, query):
        """Main search with automatic fallback and multiple sources"""
        print(f"\n{'='*60}")
        print(f"🔎 Multi-Source Search: {query}")
        print(f"{'='*60}")
        
        all_content = []
        sources_used = []
        
        # Try comparison first if it's a vs query
        if ' vs ' in query.lower() or ' versus ' in query.lower():
            text, source = self.fetch_comparison_content(query)
            if text and len(text) > 200:
                all_content.append(text)
                sources_used.append(source)
                combined_text = text
                segments = segment_into_topics(combined_text, query)
                print(f"  ✅ {len(segments)} segments from comparison")
                print(f"{'='*60}\n")
                return segments, sources_used
        
        # Try sources in priority order
        fetch_methods = [
            self.fetch_from_wikipedia,
            self.fetch_from_simple_wikipedia,
            self.fetch_from_duckduckgo,
            self.fetch_from_web_search
        ]
        
        for fetch_method in fetch_methods:
            try:
                text, source = fetch_method(query)
                if text and len(text) > 100:
                    all_content.append(text)
                    sources_used.append(source)
                    
                    # Stop if we have substantial content
                    if len(text) > 500:
                        break
                
                time.sleep(0.5)  # Rate limiting
            except Exception as e:
                print(f"  ⚠️ Error: {e}")
                continue
        
        # If no content found, use fallback
        if not all_content:
            print(f"  ⚠️ No content from primary sources, using fallback...")
            text, source = self.fetch_from_fallback_sources(query)
            if text:
                all_content.append(text)
                sources_used.append(source)
        
        if not all_content:
            print(f"  ❌ No content found")
            return [], []
        
        combined_text = "\n\n".join(all_content)
        segments = segment_into_topics(combined_text, query)
        
        print(f"  ✅ {len(segments)} segments from {len(sources_used)} sources")
        print(f"{'='*60}\n")
        
        return segments, sources_used

# Initialize multi-source learner
multi_learner = MultiSourceLearner()

# ==========================================
# FILE PROCESSING
# ==========================================
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pdf(filepath):
    text = ""
    try:
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n\n"
    except Exception as e:
        print(f"  ❌ PDF error: {e}")
    return text.strip()

def extract_text_from_ppt(filepath):
    text = ""
    try:
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    except Exception as e:
        print(f"  ❌ PPT error: {e}")
    return text

def extract_text_from_txt(filepath):
    try:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        print(f"  ❌ TXT error: {e}")
        return ""

# ==========================================
# TEXT PROCESSING
# ==========================================
def clean_text(text):
    if not text:
        return ""
    # Remove wiki markup
    text = re.sub(r'={2,}[^=]*={2,}', '', text)
    text = re.sub(r'\[.*?\]', '', text)
    # Remove excessive newlines
    text = re.sub(r'\n{3,}', '\n\n', text)
    # Remove excessive spaces
    text = re.sub(r' {2,}', ' ', text)
    # Remove non-printable characters but keep basic punctuation and common symbols
    text = re.sub(r'[^\x20-\x7E\n\r\t]', '', text)
    return text.strip()

def extract_complete_sentences(text):
    if not text:
        return []
    
    try:
        sentences = sent_tokenize(text)
    except:
        sentences = [s.strip() + '.' for s in text.split('.') if s.strip()]
    
    complete = []
    for sent in sentences:
        sent = sent.strip()
        if not sent:
            continue
        if sent[-1] not in '.!?':
            sent += '.'
        word_count = len(sent.split())
        if 3 <= word_count <= 100:
            complete.append(sent)
    
    return complete

def segment_into_topics(text, main_topic):
    """Segment text into digestible topics"""
    if not text or len(text.strip()) < 20:
        # Return a basic segment if text is too short
        return [{
            'title': main_topic,
            'content': [text] if text else [f"Learning about {main_topic}"],
            'key_points': [f"Understanding {main_topic}"]
        }]
    
    sentences = extract_complete_sentences(text)
    
    if not sentences or len(sentences) == 0:
        # If no sentences extracted, split by paragraphs
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        if paragraphs:
            sentences = paragraphs
        else:
            # Last resort: split by newlines
            sentences = [s.strip() for s in text.split('\n') if s.strip() and len(s.strip()) > 20]
    
    if not sentences:
        # Absolute fallback
        return [{
            'title': main_topic,
            'content': [text[:500]] if len(text) > 500 else [text],
            'key_points': [f"Key concepts about {main_topic}"]
        }]
    
    if len(sentences) < 3:
        return [{
            'title': main_topic,
            'content': sentences,
            'key_points': sentences
        }]
    
    segments = []
    segment_size = 5
    
    for i in range(0, len(sentences), segment_size):
        batch = sentences[i:i+segment_size+2]
        if batch:
            segments.append({
                'title': f"{main_topic} - Part {len(segments)+1}",
                'content': batch,
                'key_points': batch[:5]
            })
    
    return segments if segments else [{
        'title': main_topic,
        'content': sentences[:10],
        'key_points': sentences[:5]
    }]

def generate_quiz_for_segment(segment):
    """Generate proper MCQ questions from segment content"""
    sentences = segment.get('content', [])
    if not sentences or len(sentences) < 2:
        return []
    
    questions = []
    
    # Generate different types of questions
    for idx, sent in enumerate(sentences[:3]):
        if not sent or not isinstance(sent, str):
            continue
            
        words = sent.split()
        if len(words) < 5:
            continue
            
        # Extract key information
        important_words = [w.strip('.,!?;:') for w in words if len(w) > 4 and w[0].isupper()]
        
        if important_words and len(important_words) > 0:
            # Type 1: Fill in the blank style
            key_word = important_words[0]
            question_text = sent.replace(key_word, "______", 1)
            
            # Generate distractors (wrong answers)
            other_words = [w for w in important_words if w != key_word]
            distractors = other_words[:2] if len(other_words) >= 2 else []
            
            # Add generic distractors if needed
            generic_distractors = ["None of these", "All of the above", "Cannot be determined"]
            while len(distractors) < 3:
                if generic_distractors:
                    distractors.append(generic_distractors.pop(0))
                else:
                    distractors.append(f"Not {key_word}")
            
            options = [key_word] + distractors[:3]
            random.shuffle(options)
            
            questions.append({
                'question': f"Fill in the blank: {question_text}",
                'options': options,
                'answer': key_word,
                'type': 'Multiple Choice'
            })
        
        elif len(words) > 10:
            # Type 2: True/False based on facts
            questions.append({
                'question': f"True or False: {sent}",
                'options': ['True', 'False'],
                'answer': 'True',
                'type': 'True or False'
            })
    
    # Type 3: Comprehension question
    if len(sentences) >= 3:
        topic_words = segment.get('title', 'this topic').split('-')[0].strip()
        questions.append({
            'question': f"What is the main topic discussed in this section?",
            'options': [
                topic_words,
                "General Information",
                "Historical Facts",
                "Scientific Theory"
            ],
            'answer': topic_words,
            'type': 'Multiple Choice'
        })
    
    # Ensure we have at least one question
    if not questions and sentences:
        questions.append({
            'question': f"What is discussed in this section?",
            'options': [
                segment.get('title', 'Main Topic'),
                "Unrelated Topic",
                "General Knowledge",
                "None of these"
            ],
            'answer': segment.get('title', 'Main Topic'),
            'type': 'Multiple Choice'
        })
    
    return questions[:4]  # Return max 4 questions per segment

def suggest_related_topics(main_topic):
    suggestions_db = {
        'kubernetes': ['Docker', 'Pods', 'Container Orchestration'],
        'pods': ['Kubernetes', 'Docker', 'Containers'],
        'jenkins': ['CI/CD', 'DevOps', 'Git'],
        'docker': ['Kubernetes', 'Containers', 'Microservices'],
        'python': ['Django', 'Flask', 'Data Science'],
        'machine learning': ['Deep Learning', 'AI', 'Neural Networks'],
    }
    
    lower = main_topic.lower()
    for key, topics in suggestions_db.items():
        if key in lower or lower in key:
            return topics[:3]
    
    return ['Related Topics', 'Advanced Concepts', 'Applications']

# ==========================================
# AUTHENTICATION ROUTES
# ==========================================
@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        data = request.get_json()
        username = data.get('username', '').strip()
        password = data.get('password', '')
        
        if not username or not password:
            return jsonify({'error': 'Username and password required'}), 400
        
        conn = get_db_connection()
        user = conn.execute('SELECT * FROM users WHERE username = ?', (username,)).fetchone()
        conn.close()
        
        if user and check_password_hash(user['password_hash'], password):
            session['user_id'] = user['id']
            session['username'] = user['username']
            
            conn = get_db_connection()
            conn.execute('UPDATE users SET last_login = ? WHERE id = ?', 
                        (datetime.now(), user['id']))
            conn.commit()
            conn.close()
            
            return jsonify({'success': True, 'redirect': url_for('index')})
        else:
            return jsonify({'error': 'Invalid username or password'}), 401
    
    return render_template('login.html')

@app.route('/signup', methods=['GET', 'POST'])
def signup():
    if request.method == 'POST':
        data = request.get_json()
        username = data.get('username', '').strip()
        email = data.get('email', '').strip()
        password = data.get('password', '')
        full_name = data.get('full_name', '').strip()
        
        if not all([username, email, password, full_name]):
            return jsonify({'error': 'All fields are required'}), 400
        
        if len(username) < 3:
            return jsonify({'error': 'Username must be at least 3 characters'}), 400
        
        if len(password) < 6:
            return jsonify({'error': 'Password must be at least 6 characters'}), 400
        
        conn = get_db_connection()
        
        existing = conn.execute('SELECT id FROM users WHERE username = ?', (username,)).fetchone()
        if existing:
            conn.close()
            return jsonify({'error': 'Username already taken'}), 400
        
        existing = conn.execute('SELECT id FROM users WHERE email = ?', (email,)).fetchone()
        if existing:
            conn.close()
            return jsonify({'error': 'Email already registered'}), 400
        
        password_hash = generate_password_hash(password)
        
        try:
            cursor = conn.execute('''
                INSERT INTO users (username, email, password_hash, full_name)
                VALUES (?, ?, ?, ?)
            ''', (username, email, password_hash, full_name))
            
            user_id = cursor.lastrowid
            
            conn.execute('''
                INSERT INTO user_stats (user_id, streak, last_activity, total_quizzes, correct_answers, total_questions, total_study_time)
                VALUES (?, 0, NULL, 0, 0, 0, 0)
            ''', (user_id,))
            
            conn.commit()
            conn.close()
            
            session['user_id'] = user_id
            session['username'] = username
            
            return jsonify({'success': True, 'redirect': url_for('index')})
        
        except Exception as e:
            conn.close()
            return jsonify({'error': 'Registration failed'}), 500
    
    return render_template('signup.html')

@app.route('/logout')
def logout():
    session.clear()
    return redirect(url_for('login'))

# ==========================================
# MOTIVATIONAL QUOTES
# ==========================================
MOTIVATIONAL_QUOTES = [
    {
        "quote": "You've got this! Every expert was once a beginner.",
        "emoji": "🌟"
    },
    {
        "quote": "Small progress is still progress. Keep going!",
        "emoji": "🚀"
    },
    {
        "quote": "Believe in yourself. You're capable of amazing things!",
        "emoji": "💪"
    },
    {
        "quote": "Learning is a journey, not a race. Take your time!",
        "emoji": "🎯"
    },
    {
        "quote": "Your future self will thank you for studying today!",
        "emoji": "✨"
    },
    {
        "quote": "Mistakes are proof that you're trying. Keep learning!",
        "emoji": "🌈"
    },
    {
        "quote": "You're doing better than you think. Stay positive!",
        "emoji": "😊"
    },
    {
        "quote": "Every study session brings you closer to your goals!",
        "emoji": "🎓"
    },
    {
        "quote": "Take a deep breath. You're exactly where you need to be.",
        "emoji": "🌸"
    },
    {
        "quote": "Your dedication today creates your success tomorrow!",
        "emoji": "🏆"
    },
    {
        "quote": "Learning something new? That's courage in action!",
        "emoji": "🦁"
    },
    {
        "quote": "Progress, not perfection. You're on the right path!",
        "emoji": "🛤️"
    },
    {
        "quote": "Take breaks when needed. Rest is part of growth!",
        "emoji": "☕"
    },
    {
        "quote": "You're building your future, one lesson at a time!",
        "emoji": "🏗️"
    },
    {
        "quote": "Celebrate small wins! They lead to big achievements!",
        "emoji": "🎉"
    },
    {
        "quote": "Your effort matters more than you know. Keep it up!",
        "emoji": "💫"
    },
    {
        "quote": "Struggling? That means you're growing. Don't give up!",
        "emoji": "🌱"
    },
    {
        "quote": "You're not alone in this journey. We're here to help!",
        "emoji": "🤝"
    },
    {
        "quote": "Today's study session is tomorrow's confidence!",
        "emoji": "🔥"
    },
    {
        "quote": "Be proud of yourself for showing up and trying!",
        "emoji": "⭐"
    },
    {
        "quote": "Knowledge is power, and you're getting stronger!",
        "emoji": "💡"
    },
    {
        "quote": "One step at a time. You're making great progress!",
        "emoji": "👣"
    },
    {
        "quote": "Your brain is amazing! Feed it with curiosity!",
        "emoji": "🧠"
    },
    {
        "quote": "Consistency beats intensity. Keep showing up!",
        "emoji": "📚"
    },
    {
        "quote": "You're investing in yourself. That's the best investment!",
        "emoji": "💎"
    }
]

def get_random_quote():
    """Get a random motivational quote"""
    return random.choice(MOTIVATIONAL_QUOTES)

# ==========================================
# MAIN APP ROUTES
# ==========================================
@app.route('/')
@login_required
def index():
    user_id = session.get('user_id')
    user = get_user_by_id(user_id)
    
    # Update streak when user visits
    current_streak = update_streak(user_id)
    
    # Get a random motivational quote
    daily_quote = get_random_quote()
    
    user_data = {
        'username': user['username'],
        'full_name': user['full_name'],
        'streak': current_streak
    }
    
    return render_template('index_modern.html', user=user_data, quote=daily_quote)

@app.route('/generate', methods=['POST'])
@login_required
def generate():
    try:
        user_id = session.get('user_id')
        
        # Handle both JSON and form data
        if request.is_json:
            data = request.get_json()
            query = data.get('query', '').strip()
            user_notes = data.get('user_notes', '').strip()
            quiz_mode = data.get('quiz_mode', 'enabled')
            files = []
        else:
            query = request.form.get('query', '').strip()
            user_notes = request.form.get('user_notes', '').strip()
            quiz_mode = request.form.get('quiz_mode', 'enabled')
            files = request.files.getlist('files')
        
        print("\n" + "="*60)
        print(f"🚀 Query: {query}")
        print("="*60)
        
        if not query:
            return jsonify({'error': 'Please enter a topic'}), 400
        
        streak = update_streak(user_id)
        
        # Process uploaded files
        uploaded_content = ""
        if files:
            for file in files:
                if file and file.filename and allowed_file(file.filename):
                    try:
                        filename = secure_filename(file.filename)
                        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                        file.save(filepath)
                        
                        ext = filename.rsplit('.', 1)[1].lower()
                        
                        if ext == 'pdf':
                            uploaded_content += extract_text_from_pdf(filepath)
                        elif ext in ['ppt', 'pptx']:
                            uploaded_content += extract_text_from_ppt(filepath)
                        elif ext == 'txt':
                            uploaded_content += extract_text_from_txt(filepath)
                        
                        # Clean up file
                        try:
                            os.remove(filepath)
                        except:
                            pass
                    except Exception as e:
                        print(f"  ⚠️ Error processing file {file.filename}: {e}")
        
        segments = []
        sources = []
        
        # Priority: 1. Uploaded files, 2. User notes, 3. Multi-source search
        if uploaded_content and len(uploaded_content.strip()) > 50:
            print("✅ Using uploaded files")
            cleaned = clean_text(uploaded_content)
            segments = segment_into_topics(cleaned, query)
            sources = ['Your uploaded files']
        elif user_notes and len(user_notes.strip()) > 50:
            print("✅ Using user notes")
            segments = segment_into_topics(user_notes, query)
            sources = ['Your notes']
        else:
            print("✅ Using multi-source search")
            segments, sources = multi_learner.search_and_learn(query)
        
        if not segments:
            # Try fallback one more time
            print("⚠️ No segments from primary search, trying fallback...")
            fallback_text, fallback_source = multi_learner.fetch_from_fallback_sources(query)
            if fallback_text:
                segments = segment_into_topics(fallback_text, query)
                sources = [fallback_source]
        
        if not segments:
            print("❌ Still no segments after fallback")
            return jsonify({'error': f'Could not generate content for "{query}". Please try uploading your own notes or try a different topic.'}), 404
        
        # Generate quizzes with proper MCQs - ONE QUIZ PER SEGMENT
        for idx, segment in enumerate(segments):
            segment['segment_index'] = idx  # Add segment index
            segment['quiz'] = generate_quiz_for_segment(segment) if quiz_mode == 'enabled' else []
        
        suggestions = suggest_related_topics(query)
        
        # Generate a unique session ID for this quiz session
        import uuid
        quiz_session_id = str(uuid.uuid4())
        
        # Store quiz data in database instead of session (avoids cookie size limits)
        conn = get_db_connection()
        quiz_data_json = json.dumps({
            'segments': segments,
            'topic': query,
            'quiz_mode': quiz_mode,
            'sources': sources
        })
        conn.execute('''
            INSERT INTO temp_quiz_sessions (session_id, user_id, quiz_data, topic)
            VALUES (?, ?, ?, ?)
        ''', (quiz_session_id, user_id, quiz_data_json, query))
        conn.commit()
        conn.close()
        
        # Store only the session ID in Flask session (small footprint)
        session['quiz_session_id'] = quiz_session_id
        session['current_topic'] = query
        
        print(f"✅ Stored quiz session: {quiz_session_id}")
        print(f"✅ Returning {len(segments)} segments from {len(sources)} sources")
        print("="*60 + "\n")
        
        return jsonify({
            'success': True,
            'topic': query,
            'segments': segments,
            'sources': sources,
            'quiz_mode': quiz_mode,
            'suggestions': suggestions,
            'streak': streak,
            'quiz_session_id': quiz_session_id
        })
        
    except Exception as e:
        print(f"❌ ERROR: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': f'Server error: {str(e)}'}), 500

@app.route('/submit_quiz', methods=['POST'])
@login_required
def submit_quiz():
    try:
        user_id = session.get('user_id')
        data = request.get_json()
        
        if not data:
            return jsonify({'error': 'No data provided'}), 400
        
        segment_idx = data.get('segment_index', 0)
        answers = data.get('answers', {})
        study_duration = data.get('study_duration', 0)  # Time in seconds
        
        print(f"\n📝 Quiz Submission:")
        print(f"  User ID: {user_id}")
        print(f"  Segment: {segment_idx}")
        print(f"  Answers: {len(answers)}")
        print(f"  Study Duration: {study_duration} seconds")
        
        # Retrieve quiz data from database instead of session
        quiz_session_id = session.get('quiz_session_id') or data.get('quiz_session_id')
        
        if not quiz_session_id:
            print("  ❌ No quiz_session_id found")
            return jsonify({'error': 'No active study session. Please regenerate the content and try again.'}), 400
        
        conn = get_db_connection()
        quiz_session = conn.execute(
            'SELECT quiz_data, topic FROM temp_quiz_sessions WHERE session_id = ? AND user_id = ?',
            (quiz_session_id, user_id)
        ).fetchone()
        
        if not quiz_session:
            conn.close()
            print("  ❌ Quiz session not found in database")
            return jsonify({'error': 'Quiz session expired. Please regenerate the content and try again.'}), 400
        
        study_data = json.loads(quiz_session['quiz_data'])
        segments = study_data.get('segments', [])
        
        print(f"  Retrieved quiz session: {quiz_session_id}")
        print(f"  Segments count: {len(segments)}")
        
        if segment_idx >= len(segments):
            return jsonify({'error': 'Invalid segment'}), 400
        
        questions = segments[segment_idx].get('quiz', [])
        
        if not questions:
            conn.close()
            return jsonify({'error': 'No quiz questions available'}), 400
        
        # Calculate score
        score = 0
        for i, q in enumerate(questions):
            user_answer = answers.get(str(i))
            if user_answer and user_answer == q.get('answer'):
                score += 1
        
        # Get or create user_stats
        stats = conn.execute('SELECT * FROM user_stats WHERE user_id = ?', (user_id,)).fetchone()
        
        if not stats:
            # Create user_stats if it doesn't exist
            conn.execute('''
                INSERT INTO user_stats (user_id, streak, last_activity, total_quizzes, correct_answers, total_questions, total_study_time)
                VALUES (?, 0, NULL, 0, 0, 0, 0)
            ''', (user_id,))
            conn.commit()
            stats = conn.execute('SELECT * FROM user_stats WHERE user_id = ?', (user_id,)).fetchone()
        
        new_total = stats['total_quizzes'] + 1
        new_correct = stats['correct_answers'] + score
        new_questions = stats['total_questions'] + len(questions)
        
        # Add study time (convert seconds to minutes)
        study_minutes = int(study_duration / 60) if study_duration > 0 else 5  # Default 5 min if not tracked
        current_study_time = stats['total_study_time'] if stats['total_study_time'] is not None else 0
        new_study_time = current_study_time + study_minutes
        
        print(f"  Study Duration (seconds): {study_duration}")
        print(f"  Study Minutes: {study_minutes}")
        print(f"  Old Total Time: {current_study_time}")
        print(f"  New Total Time: {new_study_time}")
        
        conn.execute('''UPDATE user_stats 
                       SET total_quizzes = ?, correct_answers = ?, total_questions = ?, total_study_time = ?
                       WHERE user_id = ?''',
                    (new_total, new_correct, new_questions, new_study_time, user_id))
        
        # Store quiz data with questions and answers
        quiz_data = {
            'questions': [],
            'user_answers': [],
            'correct_answers': [],
            'results': []
        }
        
        for i, q in enumerate(questions):
            user_answer = answers.get(str(i), 'Not answered')
            correct_answer = q.get('answer', '')
            is_correct = user_answer == correct_answer
            
            quiz_data['questions'].append(q.get('question', ''))
            quiz_data['user_answers'].append(user_answer)
            quiz_data['correct_answers'].append(correct_answer)
            quiz_data['results'].append(is_correct)
        
        quiz_data_json = json.dumps(quiz_data)
        
        conn.execute('INSERT INTO study_history (user_id, topic, score, difficulty, study_duration, quiz_data) VALUES (?, ?, ?, ?, ?, ?)',
                    (user_id, study_data.get('topic', 'Unknown'), f"{score}/{len(questions)}", 'medium', study_minutes, quiz_data_json))
        
        conn.commit()
        conn.close()
        
        # Update streak
        update_streak(user_id)
        
        percentage = round((score / len(questions)) * 100, 1) if questions else 0
        
        print(f"  ✅ Quiz submitted successfully")
        print(f"  Score: {score}/{len(questions)} ({percentage}%)")
        
        # Prepare detailed results for each question
        detailed_results = []
        for i, q in enumerate(questions):
            user_answer = answers.get(str(i), 'Not answered')
            correct_answer = q.get('answer', '')
            is_correct = user_answer == correct_answer
            
            detailed_results.append({
                'question': q.get('question', ''),
                'options': q.get('options', []),
                'user_answer': user_answer,
                'correct_answer': correct_answer,
                'is_correct': is_correct
            })
        
        return jsonify({
            'success': True,
            'score': score,
            'total': len(questions),
            'percentage': percentage,
            'study_time_added': study_minutes,
            'detailed_results': detailed_results
        })
        
    except Exception as e:
        print(f"❌ Error in submit_quiz: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500

@app.route('/profile')
@login_required
def profile():
    user_id = session.get('user_id')
    user = get_user_by_id(user_id)
    stats = get_user_stats(user_id)
    
    # Update streak when visiting profile
    current_streak = update_streak(user_id)
    
    # Calculate overall score
    overall_score = 0
    if stats and stats['total_questions'] > 0:
        overall_score = round((stats['correct_answers'] / stats['total_questions']) * 100, 1)
    
    # Get study history
    conn = get_db_connection()
    history = conn.execute('''
        SELECT * FROM study_history 
        WHERE user_id = ? 
        ORDER BY timestamp DESC 
        LIMIT 10
    ''', (user_id,)).fetchall()
    conn.close()
    
    history_list = [dict(row) for row in history]
    
    # 12 Achievements
    achievements = [
        {'id': '3day', 'name': '🔥 3-Day Streak', 'desc': 'Study 3 days in a row', 'unlocked': current_streak >= 3},
        {'id': '7day', 'name': '⭐ Week Warrior', 'desc': 'Study 7 days in a row', 'unlocked': current_streak >= 7},
        {'id': '30day', 'name': '🏆 Month Master', 'desc': 'Study 30 days in a row', 'unlocked': current_streak >= 30},
        {'id': 'quiz10', 'name': '📝 Quiz Novice', 'desc': 'Complete 10 quizzes', 'unlocked': stats['total_quizzes'] >= 10 if stats else False},
        {'id': 'quiz50', 'name': '🎓 Quiz Expert', 'desc': 'Complete 50 quizzes', 'unlocked': stats['total_quizzes'] >= 50 if stats else False},
        {'id': 'quiz100', 'name': '👑 Quiz Legend', 'desc': 'Complete 100 quizzes', 'unlocked': stats['total_quizzes'] >= 100 if stats else False},
        {'id': 'accuracy70', 'name': '🎯 Good Aim', 'desc': 'Achieve 70%+ accuracy', 'unlocked': overall_score >= 70},
        {'id': 'accuracy80', 'name': '🎯 Sharp Shooter', 'desc': 'Achieve 80%+ accuracy', 'unlocked': overall_score >= 80},
        {'id': 'accuracy90', 'name': '🎯 Sniper', 'desc': 'Achieve 90%+ accuracy', 'unlocked': overall_score >= 90},
        {'id': 'perfect', 'name': '💯 Perfectionist', 'desc': 'Score 100% on a quiz', 'unlocked': overall_score == 100},
        {'id': 'early_bird', 'name': '🌅 Early Bird', 'desc': 'Study before 8 AM', 'unlocked': False},
        {'id': 'night_owl', 'name': '🦉 Night Owl', 'desc': 'Study after 10 PM', 'unlocked': False}
    ]

    calendar = [{'date': (datetime.now() - timedelta(days=29-i)).date().isoformat(), 'has_activity': False} for i in range(30)]

    user_data = {
        'username': user['username'],
        'full_name': user['full_name'],
        'email': user['email'],
        'bio': user.get('bio') or 'No bio yet',
        'created_at': user['created_at'],
        'streak': current_streak,
        'total_quizzes': stats['total_quizzes'] if stats else 0,
        'correct_answers': stats['correct_answers'] if stats else 0,
        'total_questions': stats['total_questions'] if stats else 0,
        'total_study_time': stats['total_study_time'] if stats else 0,
        'overall_score': overall_score
    }

    return render_template('profile.html', user=user_data, achievements=achievements, calendar=calendar, history=history_list)

@app.route('/scheduler')
@login_required
def scheduler():
    user_id = session.get('user_id')
    user = get_user_by_id(user_id)
    stats = get_user_stats(user_id)
    
    # Update streak when visiting scheduler
    current_streak = update_streak(user_id)
    
    user_data = {
        'username': user['username'],
        'full_name': user['full_name'],
        'streak': current_streak
    }
    return render_template('scheduler.html', user=user_data)

# ==========================================
# AI SCHEDULER API ROUTES
# ==========================================

@app.route('/api/generate_schedule', methods=['POST'])
@login_required
def api_generate_schedule():
    """AI-powered schedule generation"""
    try:
        user_id = session.get('user_id')
        data = request.get_json()
        
        subject = data.get('subject', '').strip()
        syllabus = data.get('syllabus', '').strip()
        total_days = int(data.get('total_days', 30))
        start_date = data.get('start_date')
        
        if not subject or not syllabus:
            return jsonify({'error': 'Subject and syllabus are required'}), 400
        
        if total_days < 1 or total_days > 365:
            return jsonify({'error': 'Total days must be between 1 and 365'}), 400
        
        # Generate schedule using AI
        schedule = AIScheduler.generate_schedule(syllabus, total_days, start_date)
        
        if schedule is None or len(schedule) == 0:
            return jsonify({'error': 'Could not parse syllabus. Please check the format and ensure it contains topics.'}), 400
        
        # Add subject to each item
        for item in schedule:
            item['subject'] = subject
        
        return jsonify({
            'success': True,
            'schedule': schedule,
            'total_topics': len(schedule),
            'total_hours': sum(item['estimated_hours'] for item in schedule)
        })
        
    except Exception as e:
        print(f"Error generating schedule: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500

@app.route('/api/save_bulk_schedule', methods=['POST'])
@login_required
def api_save_bulk_schedule():
    """Save AI-generated schedule to database"""
    try:
        user_id = session.get('user_id')
        data = request.get_json()
        schedule = data.get('schedule', [])
        
        if not schedule:
            return jsonify({'error': 'No schedule data provided'}), 400
        
        conn = get_db_connection()
        
        for item in schedule:
            conn.execute('''
                INSERT INTO schedules 
                (user_id, subject, unit_name, topic_name, difficulty, estimated_hours, 
                 scheduled_date, start_time, end_time, is_auto_generated, status)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, 'pending')
            ''', (
                user_id,
                item.get('subject', ''),
                item.get('unit_name', ''),
                item.get('topic_name', ''),
                item.get('difficulty', 'medium'),
                item.get('estimated_hours', 2),
                item.get('scheduled_date', ''),
                item.get('start_time', '09:00'),
                item.get('end_time', '11:00'),
                item.get('is_auto_generated', 1)
            ))
        
        conn.commit()
        conn.close()
        
        return jsonify({'success': True, 'saved_count': len(schedule)})
        
    except Exception as e:
        print(f"Error saving schedule: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/add_schedule_item', methods=['POST'])
@login_required
def api_add_schedule_item():
    """Add a single manual schedule item"""
    try:
        user_id = session.get('user_id')
        data = request.get_json()
        
        subject = data.get('subject', '').strip()
        topic_name = data.get('topic_name', '').strip()
        scheduled_date = data.get('scheduled_date', '')
        
        if not all([subject, topic_name, scheduled_date]):
            return jsonify({'error': 'Missing required fields'}), 400
        
        conn = get_db_connection()
        conn.execute('''
            INSERT INTO schedules 
            (user_id, subject, unit_name, topic_name, difficulty, estimated_hours, 
             scheduled_date, start_time, end_time, notes, is_auto_generated, status)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, 0, 'pending')
        ''', (
            user_id,
            subject,
            data.get('unit_name', ''),
            topic_name,
            data.get('difficulty', 'medium'),
            data.get('estimated_hours', 2),
            scheduled_date,
            data.get('start_time', '09:00'),
            data.get('end_time', '11:00'),
            data.get('notes', '')
        ))
        
        conn.commit()
        conn.close()
        
        return jsonify({'success': True})
        
    except Exception as e:
        print(f"Error adding schedule item: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/get_schedule')
@login_required
def api_get_schedule():
    """Get all schedules for current user"""
    try:
        user_id = session.get('user_id')
        
        conn = get_db_connection()
        schedules = conn.execute('''
            SELECT * FROM schedules 
            WHERE user_id = ? 
            ORDER BY scheduled_date, start_time
        ''', (user_id,)).fetchall()
        conn.close()
        
        schedule_list = [dict(row) for row in schedules]
        
        return jsonify({
            'success': True,
            'schedule': schedule_list
        })
        
    except Exception as e:
        print(f"Error getting schedule: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/complete_schedule/<int:schedule_id>', methods=['POST'])
@login_required
def api_complete_schedule(schedule_id):
    """Mark a schedule item as completed"""
    try:
        user_id = session.get('user_id')
        
        conn = get_db_connection()
        conn.execute('''
            UPDATE schedules 
            SET status = 'completed', 
                completion_percentage = 100,
                completed_at = ?
            WHERE id = ? AND user_id = ?
        ''', (datetime.now(), schedule_id, user_id))
        conn.commit()
        conn.close()
        
        # Update streak
        update_streak(user_id)
        
        return jsonify({'success': True})
        
    except Exception as e:
        print(f"Error completing schedule: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/delete_schedule/<int:schedule_id>', methods=['DELETE'])
@login_required
def api_delete_schedule(schedule_id):
    """Delete a schedule item"""
    try:
        user_id = session.get('user_id')
        
        conn = get_db_connection()
        conn.execute('DELETE FROM schedules WHERE id = ? AND user_id = ?', 
                    (schedule_id, user_id))
        conn.commit()
        conn.close()
        
        return jsonify({'success': True})
        
    except Exception as e:
        print(f"Error deleting schedule: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/export_schedule')
@login_required
def api_export_schedule():
    """Export schedule as text file"""
    try:
        user_id = session.get('user_id')
        
        conn = get_db_connection()
        schedules = conn.execute('''
            SELECT * FROM schedules 
            WHERE user_id = ? 
            ORDER BY scheduled_date, start_time
        ''', (user_id,)).fetchall()
        conn.close()
        
        # Create text export
        export_text = "MY AI-GENERATED STUDY SCHEDULE\n"
        export_text += "=" * 60 + "\n\n"
        
        # Group by date
        from collections import defaultdict
        by_date = defaultdict(list)
        
        for schedule in schedules:
            by_date[schedule['scheduled_date']].append(schedule)
        
        for date in sorted(by_date.keys()):
            date_obj = datetime.strptime(date, '%Y-%m-%d')
            formatted_date = date_obj.strftime('%A, %B %d, %Y')
            
            export_text += f"\n{formatted_date}\n"
            export_text += "-" * 60 + "\n"
            
            for schedule in by_date[date]:
                status_icon = "✓" if schedule['status'] == 'completed' else "○"
                export_text += f"{status_icon} {schedule['start_time']} - {schedule['end_time']}: "
                export_text += f"{schedule['topic_name']} ({schedule['difficulty'].upper()})\n"
                
                if schedule['unit_name']:
                    export_text += f"   Unit: {schedule['unit_name']}\n"
                
                if schedule['notes']:
                    export_text += f"   Notes: {schedule['notes']}\n"
                
                export_text += f"   Estimated Time: {schedule['estimated_hours']} hours\n"
                export_text += "\n"
        
        # Add statistics
        total = len(schedules)
        completed = sum(1 for s in schedules if s['status'] == 'completed')
        total_hours = sum(s['estimated_hours'] for s in schedules)
        
        export_text += "\n" + "=" * 60 + "\n"
        export_text += "STATISTICS\n"
        export_text += "=" * 60 + "\n"
        export_text += f"Total Topics: {total}\n"
        export_text += f"Completed: {completed}\n"
        export_text += f"Remaining: {total - completed}\n"
        export_text += f"Total Study Hours: {total_hours}\n"
        export_text += f"Progress: {round((completed/total)*100, 1) if total > 0 else 0}%\n"
        
        return Response(
            export_text,
            mimetype="text/plain",
            headers={"Content-disposition": "attachment; filename=study_schedule.txt"}
        )
        
    except Exception as e:
        print(f"Error exporting schedule: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/download_summary', methods=['POST'])
@login_required
def api_download_summary():
    """Download summary as text file"""
    try:
        data = request.get_json()
        topic = data.get('topic', 'Summary')
        segments = data.get('segments', [])
        sources = data.get('sources', [])
        
        if not segments:
            return jsonify({'error': 'No summary data provided'}), 400
        
        # Create formatted text export
        export_text = f"STUDY SUMMARY: {topic}\n"
        export_text += "=" * 80 + "\n\n"
        
        # Add sources
        if sources:
            export_text += "SOURCES:\n"
            for idx, source in enumerate(sources, 1):
                export_text += f"{idx}. {source}\n"
            export_text += "\n" + "-" * 80 + "\n\n"
        
        # Add segments
        for idx, segment in enumerate(segments, 1):
            export_text += f"\n{idx}. {segment.get('title', f'Topic {idx}')}\n"
            export_text += "-" * 80 + "\n\n"
            
            # Add content
            content = segment.get('content', [])
            if isinstance(content, list):
                for item in content:
                    export_text += f"• {item}\n"
            else:
                export_text += f"{content}\n"
            
            export_text += "\n"
            
            # Add key points if available
            key_points = segment.get('key_points', [])
            if key_points:
                export_text += "KEY POINTS:\n"
                for point in key_points:
                    export_text += f"  ✓ {point}\n"
                export_text += "\n"
        
        # Add footer
        export_text += "\n" + "=" * 80 + "\n"
        export_text += f"Generated by StudyPal AI\n"
        export_text += f"Date: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
        
        # Create safe filename
        safe_topic = "".join(c for c in topic if c.isalnum() or c in (' ', '-', '_')).strip()
        filename = f"summary_{safe_topic[:50]}.txt"
        
        return Response(
            export_text,
            mimetype="text/plain",
            headers={"Content-disposition": f"attachment; filename={filename}"}
        )
        
    except Exception as e:
        print(f"Error downloading summary: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/get_streak')
@login_required
def api_get_streak():
    """Get current user streak"""
    try:
        user_id = session.get('user_id')
        stats = get_user_stats(user_id)
        
        return jsonify({
            'success': True,
            'streak': stats['streak'] if stats else 0
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/get_quiz_history/<int:history_id>')
@login_required
def api_get_quiz_history(history_id):
    """Get detailed quiz history including questions and answers"""
    try:
        user_id = session.get('user_id')
        
        conn = get_db_connection()
        history = conn.execute('''
            SELECT * FROM study_history 
            WHERE id = ? AND user_id = ?
        ''', (history_id, user_id)).fetchone()
        conn.close()
        
        if not history:
            return jsonify({'error': 'Quiz history not found'}), 404
        
        quiz_data = None
        if history['quiz_data']:
            try:
                quiz_data = json.loads(history['quiz_data'])
            except:
                quiz_data = None
        
        return jsonify({
            'success': True,
            'topic': history['topic'],
            'score': history['score'],
            'timestamp': history['timestamp'],
            'study_duration': history['study_duration'],
            'quiz_data': quiz_data
        })
        
    except Exception as e:
        print(f"Error getting quiz history: {e}")
        return jsonify({'error': str(e)}), 500
# ==========================================
# RUN APP
# ==========================================
if __name__ == '__main__':
    print("\n" + "="*60)
    print("🎓 AI Study Pal v3.0 - MULTI-SOURCE EDITION")
    print("="*60)
    print("🌐 Open: http://127.0.0.1:5000")
    print("🌍 Sources: Wikipedia, Simple Wiki, DuckDuckGo")
    print("📅 Features: Scheduler, Achievements, Progressive Learning")
    print("="*60 + "\n")
    app.run(debug=True, host='0.0.0.0', port=5000)